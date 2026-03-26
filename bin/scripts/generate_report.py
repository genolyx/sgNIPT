#!/usr/bin/env python3
"""
generate_report.py - Final Report Generator for Single Gene NIPT Pipeline

Generates two types of reports:
  1. JSON + HTML internal report (pipeline output)
  2. PPTX template-based PDF report (daemon-triggered, client-facing)

The PPTX report follows the same pattern as the sWGS NIPT daemon:
  - Load PPTX template with placeholders (<<key>>)
  - Replace placeholders with analysis results
  - Convert PPTX to PDF via LibreOffice
  - Merge multi-language PDFs if needed
"""

import argparse
import json
import logging
import os
import re
import shutil
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════════════
# Part 1: Pipeline Internal Report (JSON + HTML)
# ═══════════════════════════════════════════════════════════════════════════════

def load_json(path: str) -> Optional[Dict[str, Any]]:
    """Safely load a JSON file."""
    p = Path(path)
    if not p.exists():
        logger.warning("File not found: %s", path)
        return None
    with open(p, "r") as fh:
        return json.load(fh)


def aggregate_results(
    sample_id: str,
    fastq_qc_json: Optional[str],
    bam_qc_json: Optional[str],
    fetal_fraction_json: Optional[str],
    variant_report_json: Optional[str],
    multiqc_data_json: Optional[str],
) -> Dict[str, Any]:
    """Aggregate all pipeline results into a single report."""
    logger.info("Aggregating results for sample: %s", sample_id)

    report = {
        "report_metadata": {
            "sample_id": sample_id,
            "pipeline": "Single Gene NIPT Pipeline",
            "version": "1.0.0",
            "generated_at": datetime.now().isoformat(),
            "report_type": "single_gene_nipt",
        },
        "overall_status": "PASS",
        "status_flags": [],
        "fastq_qc": None,
        "bam_qc": None,
        "fetal_fraction": None,
        "variant_analysis": None,
    }

    # Load FASTQ QC
    if fastq_qc_json:
        data = load_json(fastq_qc_json)
        if data:
            report["fastq_qc"] = data
            if not data.get("qc_evaluation", {}).get("overall_qc_pass", True):
                report["status_flags"].append("FASTQ_QC_FAIL")
                report["overall_status"] = "WARNING"

    # Load BAM QC
    if bam_qc_json:
        data = load_json(bam_qc_json)
        if data:
            report["bam_qc"] = data
            if not data.get("qc_evaluation", {}).get("overall_qc_pass", True):
                report["status_flags"].append("BAM_QC_FAIL")
                report["overall_status"] = "WARNING"

    # Load Fetal Fraction
    if fetal_fraction_json:
        data = load_json(fetal_fraction_json)
        if data:
            report["fetal_fraction"] = data
            ff_status = data.get("status", "")
            if ff_status == "FAILED":
                report["status_flags"].append("FF_ESTIMATION_FAILED")
                report["overall_status"] = "FAIL"
            elif ff_status == "LOW_FF":
                report["status_flags"].append("LOW_FETAL_FRACTION")
                report["overall_status"] = "WARNING"
            elif ff_status == "INSUFFICIENT_DATA":
                report["status_flags"].append("INSUFFICIENT_SNP_DATA")
                report["overall_status"] = "WARNING"

    # Load Variant Analysis
    if variant_report_json:
        data = load_json(variant_report_json)
        if data:
            report["variant_analysis"] = data

    # Final status determination
    if report["overall_status"] == "FAIL":
        logger.warning("Overall status: FAIL for sample %s", sample_id)
    elif report["status_flags"]:
        logger.warning(
            "Overall status: WARNING for sample %s (flags: %s)",
            sample_id,
            ", ".join(report["status_flags"]),
        )
    else:
        logger.info("Overall status: PASS for sample %s", sample_id)

    return report


# ═══════════════════════════════════════════════════════════════════════════════
# Part 2: PPTX Template-Based PDF Report (daemon-compatible)
# ═══════════════════════════════════════════════════════════════════════════════

def get_template_files(template_dir: str) -> Dict[str, str]:
    """
    Scan template directory and build a mapping of template_key → filename.
    Template files follow the naming convention: {TemplateKey}_{LANG}.pptx
    e.g., "SgNIPT_Basic_EN.pptx" → key "SgNIPT_Basic_EN"
    """
    templates = {}
    template_path = Path(template_dir)

    if not template_path.exists():
        logger.warning("Template directory not found: %s", template_dir)
        return templates

    for f in template_path.glob("*.pptx"):
        if f.name.startswith("~$"):
            continue  # Skip temp files
        key = f.stem  # filename without extension
        templates[key] = f.name
        logger.debug("Found template: %s → %s", key, f.name)

    logger.info("Found %d template files in %s", len(templates), template_dir)
    return templates


def find_client_template_directory(
    base_template_dir: str, client_name: Optional[str]
) -> str:
    """
    Find client-specific template directory.
    Falls back to base directory if client-specific one doesn't exist.
    """
    if client_name:
        client_dir = os.path.join(base_template_dir, client_name)
        if os.path.isdir(client_dir):
            logger.info("Using client template directory: %s", client_dir)
            return client_dir

    logger.info("Using default template directory: %s", base_template_dir)
    return base_template_dir


def build_report_placeholders(report_json: Dict[str, Any]) -> Dict[str, str]:
    """
    Convert analysis results into placeholder key-value pairs for PPTX template.
    Placeholder format: <<key>> in the template.
    """
    placeholders = {}

    # Order metadata
    placeholders["Order ID"] = str(report_json.get("order_id", ""))
    placeholders["Sample ID"] = str(report_json.get("sample_id", ""))
    placeholders["Report Date"] = datetime.now().strftime("%Y-%m-%d")
    placeholders["Pipeline Version"] = str(report_json.get("pipeline_version", "1.0.0"))
    placeholders["Result"] = str(report_json.get("status", ""))

    # Patient info (from daemon's make_report_json)
    placeholders["Patient Name"] = str(report_json.get("Patient Name", ""))
    placeholders["Patient DOB"] = str(report_json.get("Patient DOB", ""))
    placeholders["Maternal Age"] = str(report_json.get("Maternal Age", ""))
    placeholders["Gestational Age"] = str(report_json.get("Gestational Age", ""))
    placeholders["Sample Collection Date"] = str(report_json.get("Sample Collection Date", ""))
    placeholders["Hospital"] = str(report_json.get("Hospital", ""))
    placeholders["Doctor"] = str(report_json.get("Doctor", ""))

    # Fetal Fraction
    ff_data = report_json.get("fetal_fraction", {})
    if ff_data:
        ff_value = ff_data.get("primary_ff")
        if ff_value is not None:
            placeholders["FF"] = f"{ff_value:.1%}" if isinstance(ff_value, float) else str(ff_value)
        else:
            placeholders["FF"] = "N/A"
        placeholders["FF Method"] = str(ff_data.get("primary_method", "SNP-based"))
        placeholders["FF Status"] = str(ff_data.get("status", ""))
        placeholders["Informative SNPs"] = str(ff_data.get("num_informative_snps", ""))
    else:
        placeholders["FF"] = "N/A"
        placeholders["FF Method"] = "N/A"
        placeholders["FF Status"] = "N/A"
        placeholders["Informative SNPs"] = "N/A"

    # QC Summary
    bam_qc = report_json.get("bam_qc", {})
    if bam_qc:
        placeholders["Mean Coverage"] = str(bam_qc.get("mean_coverage", "N/A"))
        placeholders["Target Coverage"] = str(bam_qc.get("target_coverage", "N/A"))
        placeholders["On-Target Rate"] = str(bam_qc.get("on_target_rate", "N/A"))
        placeholders["Mapping Rate"] = str(bam_qc.get("mapping_rate", "N/A"))

    # Variant results
    variant_data = report_json.get("variant_analysis", {})
    if variant_data:
        placeholders["Total Loci"] = str(variant_data.get("total_target_loci", ""))
        placeholders["Covered Loci"] = str(variant_data.get("covered_loci", ""))
        placeholders["Total Variants"] = str(variant_data.get("total_variants", ""))
        placeholders["Fetal Variants"] = str(variant_data.get("fetal_variants", ""))
        placeholders["Pathogenic Count"] = str(variant_data.get("pathogenic_variants", "0"))

        # Build variant detail table placeholders
        pathogenic = variant_data.get("pathogenic_details", [])
        for idx, var in enumerate(pathogenic[:10], start=1):
            prefix = f"Var{idx}"
            placeholders[f"{prefix} Gene"] = str(var.get("gene", ""))
            placeholders[f"{prefix} Position"] = str(var.get("position", ""))
            placeholders[f"{prefix} Ref"] = str(var.get("ref", ""))
            placeholders[f"{prefix} Alt"] = str(var.get("alt", ""))
            placeholders[f"{prefix} AF"] = str(var.get("allele_frequency", ""))
            placeholders[f"{prefix} Origin"] = str(var.get("origin", ""))
            placeholders[f"{prefix} Classification"] = str(var.get("classification", ""))
            placeholders[f"{prefix} Disease"] = str(var.get("disease", ""))

    return placeholders


def populate_ppt(
    template_path: str,
    output_path: str,
    placeholders: Dict[str, str],
    lang: str = "EN",
) -> None:
    """
    Replace <<placeholder>> markers in a PPTX template with actual values.
    Uses python-pptx library.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        raise

    logger.info("Populating template: %s → %s (lang=%s)", template_path, output_path, lang)

    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    modified = False

                    for key, value in placeholders.items():
                        placeholder = f"<<{key}>>"
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)
                            modified = True

                    if modified:
                        logger.debug("Replaced in slide: '%s' → '%s'", original_text, run.text)

            # Also check tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for key, value in placeholders.items():
                                    placeholder = f"<<{key}>>"
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, value)

    prs.save(output_path)
    logger.info("PPTX saved: %s", output_path)


def convert_to_pdf(pptx_path: str, output_dir: str) -> Optional[str]:
    """Convert PPTX to PDF using LibreOffice."""
    try:
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            pptx_path,
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

        if result.returncode != 0:
            logger.error("LibreOffice conversion failed: %s", result.stderr)
            return None

        pdf_path = os.path.join(
            output_dir,
            Path(pptx_path).stem + ".pdf",
        )

        if os.path.exists(pdf_path):
            logger.info("PDF generated: %s", pdf_path)
            return pdf_path
        else:
            logger.error("PDF file not found after conversion")
            return None

    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        return None
    except FileNotFoundError:
        logger.error("LibreOffice not found. Install with: apt install libreoffice")
        return None


def merge_pdfs(pdf_paths: List[str], output_path: str) -> Optional[str]:
    """Merge multiple PDFs into one."""
    try:
        from PyPDF2 import PdfMerger
    except ImportError:
        try:
            from pypdf import PdfMerger
        except ImportError:
            logger.error("PyPDF2/pypdf not installed")
            return None

    merger = PdfMerger()
    for pdf in pdf_paths:
        if os.path.exists(pdf):
            merger.append(pdf)

    merger.write(output_path)
    merger.close()

    if os.path.exists(output_path):
        logger.info("Merged PDF: %s (%d files)", output_path, len(pdf_paths))
        return output_path
    return None


def generate_from_json(
    report_json: Dict[str, Any],
    template_dir: str,
    output_dir: str,
) -> Dict[str, Any]:
    """
    Generate PDF report from JSON data using PPTX template.
    Compatible with nipt-daemon report generation pattern.

    Args:
        report_json: Report data (merged from daemon's make_report_json + analysis results)
        template_dir: Template directory path
        output_dir: Output directory path

    Returns:
        Dict with 'merged' (path) and 'individual' (list of paths)
    """
    logger.info("Generating report from JSON data")

    try:
        os.makedirs(output_dir, exist_ok=True)

        # Find client-specific template directory
        client_name = report_json.get("Client Name")
        actual_template_dir = find_client_template_directory(template_dir, client_name)

        # Get available templates
        template_files = get_template_files(actual_template_dir)
        if not template_files:
            logger.warning("No templates in %s, trying default", actual_template_dir)
            template_files = get_template_files(template_dir)
            actual_template_dir = template_dir

        if not template_files:
            raise FileNotFoundError(f"No template files found in {template_dir}")

        # Prepare placeholders
        placeholders = build_report_placeholders(report_json)

        # Language handling
        langs = report_json.get("Report Language", [])
        if isinstance(langs, str):
            langs = [langs]
        if not langs:
            langs = ["EN"]

        order_id = str(report_json.get("order_id", report_json.get("Order ID", "Report"))).strip().replace(" ", "_")
        sample_id = str(report_json.get("sample_id", report_json.get("Sample ID", "None"))).strip().replace(" ", "_")

        # Determine template key
        result_status = str(report_json.get("status", report_json.get("Result", ""))).strip().lower()
        if result_status == "no_call":
            template_key = "No_Call"
        else:
            template_key = report_json.get("TemplateKey", report_json.get("template_key", "SgNIPT_Basic"))
            template_key = re.sub(r"\s+", "_", str(template_key).strip())

        individual_pdfs = []

        for lang in [l.strip().upper() for l in langs]:
            try:
                full_key = f"{template_key}_{lang}"
                template_filename = template_files.get(full_key)

                if not template_filename:
                    logger.warning("Template '%s' not found, skipping", full_key)
                    continue

                template_path = os.path.join(actual_template_dir, template_filename)
                pptx_filename = os.path.join(output_dir, f"{order_id}_{lang}.pptx")
                pdf_filename = os.path.join(output_dir, f"{order_id}_{lang}.pdf")

                # Generate PPTX
                populate_ppt(template_path, pptx_filename, placeholders, lang)

                # Convert to PDF
                convert_to_pdf(pptx_filename, output_dir)

                if os.path.exists(pdf_filename):
                    individual_pdfs.append(pdf_filename)
                    logger.info("Generated PDF: %s", pdf_filename)

                # Clean up temp PPTX
                if os.path.exists(pptx_filename):
                    os.remove(pptx_filename)

            except Exception as e:
                logger.error("Error generating report for language %s: %s", lang, e)
                continue

        # Merge PDFs
        merged_pdf = None
        if individual_pdfs:
            merged_name = f"{order_id}_{sample_id}.pdf" if sample_id != "None" else f"{order_id}.pdf"
            merged_pdf_path = os.path.join(output_dir, merged_name)

            if len(individual_pdfs) == 1:
                shutil.copy2(individual_pdfs[0], merged_pdf_path)
                merged_pdf = merged_pdf_path
            else:
                merged_pdf = merge_pdfs(individual_pdfs, merged_pdf_path)

        result = {
            "merged": merged_pdf,
            "individual": individual_pdfs,
        }

        logger.info("Report generation completed. Generated %d individual PDFs", len(individual_pdfs))
        return result

    except Exception as e:
        logger.error("Error generating report: %s", e)
        raise


# ═══════════════════════════════════════════════════════════════════════════════
# Part 3: HTML Report (standalone, for pipeline output)
# ═══════════════════════════════════════════════════════════════════════════════

def generate_html_report(report: Dict[str, Any]) -> str:
    """Generate an HTML report from the aggregated results."""
    sample_id = report["report_metadata"]["sample_id"]
    generated_at = report["report_metadata"]["generated_at"]
    overall_status = report["overall_status"]

    status_color = {
        "PASS": "#28a745",
        "WARNING": "#ffc107",
        "FAIL": "#dc3545",
    }.get(overall_status, "#6c757d")

    html_parts = []

    # Header
    html_parts.append(f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Single Gene NIPT Report - {sample_id}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: #f5f5f5; color: #333; line-height: 1.6; }}
        .container {{ max-width: 1100px; margin: 0 auto; padding: 20px; }}
        .header {{ background: linear-gradient(135deg, #1a237e, #283593); color: white; padding: 30px; border-radius: 8px; margin-bottom: 20px; }}
        .header h1 {{ font-size: 1.8em; margin-bottom: 5px; }}
        .header .subtitle {{ opacity: 0.8; font-size: 0.95em; }}
        .status-badge {{ display: inline-block; padding: 6px 16px; border-radius: 20px; font-weight: bold; font-size: 0.9em; color: white; }}
        .card {{ background: white; border-radius: 8px; padding: 24px; margin-bottom: 16px; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
        .card h2 {{ color: #1a237e; margin-bottom: 16px; font-size: 1.3em; border-bottom: 2px solid #e8eaf6; padding-bottom: 8px; }}
        .card h3 {{ color: #283593; margin: 16px 0 8px; font-size: 1.1em; }}
        table {{ width: 100%; border-collapse: collapse; margin: 12px 0; }}
        th, td {{ padding: 10px 14px; text-align: left; border-bottom: 1px solid #e0e0e0; font-size: 0.9em; }}
        th {{ background: #e8eaf6; color: #1a237e; font-weight: 600; }}
        tr:hover {{ background: #f5f5f5; }}
        .metric-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 12px; margin: 12px 0; }}
        .metric-box {{ background: #f8f9fa; border-radius: 6px; padding: 14px; text-align: center; border: 1px solid #e0e0e0; }}
        .metric-box .value {{ font-size: 1.5em; font-weight: bold; color: #1a237e; }}
        .metric-box .label {{ font-size: 0.8em; color: #666; margin-top: 4px; }}
        .pass {{ color: #28a745; }}
        .fail {{ color: #dc3545; }}
        .warn {{ color: #ffc107; }}
        .footer {{ text-align: center; padding: 20px; color: #999; font-size: 0.85em; }}
    </style>
</head>
<body>
<div class="container">
    <div class="header">
        <h1>Single Gene NIPT Report</h1>
        <div class="subtitle">Non-Invasive Prenatal Testing - Single Gene Disorder Screening</div>
        <div style="margin-top: 15px;">
            <span><strong>Sample:</strong> {sample_id}</span> &nbsp;|&nbsp;
            <span><strong>Date:</strong> {generated_at[:10]}</span> &nbsp;|&nbsp;
            <span class="status-badge" style="background:{status_color};">{overall_status}</span>
        </div>
    </div>
""")

    # Fetal Fraction Section
    ff_data = report.get("fetal_fraction")
    if ff_data:
        ff_value = ff_data.get("primary_fetal_fraction")
        ff_method = ff_data.get("primary_method") or "N/A"
        ff_status = ff_data.get("status") or "N/A"
        ff_display = f"{ff_value:.1%}" if ff_value is not None else "N/A"
        ff_class = "pass" if ff_status == "PASS" else ("fail" if ff_status in ("FAILED", "LOW_FF") else "warn")

        snp_method = ff_data.get("methods", {}).get("snp_informative", {})
        ci = snp_method.get("confidence_interval_95", {})
        ci_text = f"{ci.get('lower', 0):.1%} - {ci.get('upper', 0):.1%}" if ci else "N/A"

        html_parts.append(f"""
    <div class="card">
        <h2>Fetal Fraction Estimation</h2>
        <div class="metric-grid">
            <div class="metric-box">
                <div class="value {ff_class}">{ff_display}</div>
                <div class="label">Fetal Fraction</div>
            </div>
            <div class="metric-box">
                <div class="value">{(ff_method or 'N/A').replace('_', ' ').title()}</div>
                <div class="label">Primary Method</div>
            </div>
            <div class="metric-box">
                <div class="value">{snp_method.get('num_informative_snps', 'N/A')}</div>
                <div class="label">Informative SNPs</div>
            </div>
            <div class="metric-box">
                <div class="value">{ci_text}</div>
                <div class="label">95% CI</div>
            </div>
        </div>
    </div>
""")

    # Variant Analysis Section
    va_data = report.get("variant_analysis")
    if va_data:
        summary = va_data.get("summary", {})
        variants = va_data.get("variants", [])

        html_parts.append(f"""
    <div class="card">
        <h2>Variant Analysis</h2>
        <div class="metric-grid">
            <div class="metric-box">
                <div class="value">{summary.get('total_target_loci', 'N/A')}</div>
                <div class="label">Target Loci</div>
            </div>
            <div class="metric-box">
                <div class="value">{summary.get('covered_loci', 'N/A')}</div>
                <div class="label">Covered Loci</div>
            </div>
            <div class="metric-box">
                <div class="value">{summary.get('total_variants_detected', 0)}</div>
                <div class="label">Variants Detected</div>
            </div>
            <div class="metric-box">
                <div class="value">{summary.get('fetal_specific_variants', 0)}</div>
                <div class="label">Fetal-Specific Variants</div>
            </div>
        </div>
""")

        if variants:
            html_parts.append("""
        <h3>Detected Variants</h3>
        <table>
            <tr>
                <th>Gene</th><th>Position</th><th>Ref/Alt</th>
                <th>AF</th><th>Origin</th><th>Classification</th>
            </tr>
""")
            for v in variants[:20]:
                cls = v.get("classification", "")
                row_style = ' style="background:#ffebee;"' if cls in ("PATHOGENIC", "LIKELY_PATHOGENIC") else ""
                html_parts.append(f"""
            <tr{row_style}>
                <td>{v.get('gene', '')}</td>
                <td>{v.get('position', '')}</td>
                <td>{v.get('ref', '')}/{v.get('alt', '')}</td>
                <td>{v.get('allele_frequency', '')}</td>
                <td>{v.get('origin', '')}</td>
                <td>{cls}</td>
            </tr>
""")
            html_parts.append("        </table>\n")

        html_parts.append("    </div>\n")

    # QC Summary Section
    fq_data = report.get("fastq_qc")
    bam_data = report.get("bam_qc")
    if fq_data or bam_data:
        html_parts.append("""
    <div class="card">
        <h2>Quality Control Summary</h2>
""")
        if fq_data:
            metrics = fq_data.get("metrics", {})
            html_parts.append(f"""
        <h3>FASTQ QC</h3>
        <table>
            <tr><th>Metric</th><th>Value</th></tr>
            <tr><td>Total Reads</td><td>{metrics.get('total_reads_before_filter', 'N/A'):,}</td></tr>
            <tr><td>Reads After Filter</td><td>{metrics.get('total_reads_after_filter', 'N/A'):,}</td></tr>
            <tr><td>Q30 Rate</td><td>{metrics.get('q30_rate_after_filter', 'N/A')}</td></tr>
            <tr><td>GC Content</td><td>{metrics.get('gc_content_after_filter', 'N/A')}</td></tr>
        </table>
""")
        if bam_data:
            metrics = bam_data.get("metrics", {})
            html_parts.append(f"""
        <h3>BAM QC</h3>
        <table>
            <tr><th>Metric</th><th>Value</th></tr>
            <tr><td>Mapped Reads</td><td>{metrics.get('mapped_reads', 'N/A')}</td></tr>
            <tr><td>Mapping Rate</td><td>{metrics.get('mapping_rate', 'N/A')}</td></tr>
            <tr><td>Mean Coverage</td><td>{metrics.get('mean_coverage', 'N/A')}</td></tr>
            <tr><td>Target Coverage</td><td>{metrics.get('target_mean_coverage', 'N/A')}</td></tr>
            <tr><td>On-Target Rate</td><td>{metrics.get('on_target_rate', 'N/A')}</td></tr>
            <tr><td>Duplicate Rate</td><td>{metrics.get('duplicate_rate', 'N/A')}</td></tr>
        </table>
""")
        html_parts.append("    </div>\n")

    # Footer
    html_parts.append(f"""
    <div class="footer">
        <p>Single Gene NIPT Pipeline v1.0.0 | Generated: {generated_at}</p>
        <p>This report is for research use only.</p>
    </div>
</div>
</body>
</html>
""")

    return "".join(html_parts)


# ═══════════════════════════════════════════════════════════════════════════════
# Main CLI
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="Generate Single Gene NIPT reports (JSON + HTML + PPTX/PDF)"
    )
    parser.add_argument("--sample_id", required=True, help="Sample ID")
    parser.add_argument("--fastq_qc", help="FASTQ QC JSON file")
    parser.add_argument("--bam_qc", help="BAM QC JSON file")
    parser.add_argument("--fetal_fraction", help="Fetal fraction JSON file")
    parser.add_argument("--variant_report", help="Variant report JSON file")
    parser.add_argument("--multiqc_data", help="MultiQC data JSON file")
    parser.add_argument("--output_dir", required=True, help="Output directory")
    parser.add_argument("--output_prefix", help="Output file prefix (default: sample_id)")

    # PPTX template options (for daemon-triggered report generation)
    parser.add_argument("--template_dir", help="PPTX template directory")
    parser.add_argument("--report_json", help="Pre-built report JSON for PPTX generation")

    args = parser.parse_args()

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = args.output_prefix or args.sample_id

    # Mode 1: PPTX template-based PDF report
    if args.report_json and args.template_dir:
        logger.info("Mode: PPTX template-based PDF report")
        report_data = load_json(args.report_json)
        if report_data:
            result = generate_from_json(report_data, args.template_dir, str(output_dir))
            logger.info("PDF report result: %s", result)
        return

    # Mode 2: Pipeline internal report (JSON + HTML)
    logger.info("Mode: Pipeline internal report (JSON + HTML)")

    report = aggregate_results(
        sample_id=args.sample_id,
        fastq_qc_json=args.fastq_qc,
        bam_qc_json=args.bam_qc,
        fetal_fraction_json=args.fetal_fraction,
        variant_report_json=args.variant_report,
        multiqc_data_json=args.multiqc_data,
    )

    # Write JSON report
    json_path = output_dir / f"{prefix}_final_report.json"
    with open(json_path, "w") as fh:
        json.dump(report, fh, indent=2, default=str)
    logger.info("JSON report written: %s", json_path)

    # Write HTML report
    html_content = generate_html_report(report)
    html_path = output_dir / f"{prefix}_final_report.html"
    with open(html_path, "w") as fh:
        fh.write(html_content)
    logger.info("HTML report written: %s", html_path)


if __name__ == "__main__":
    main()
